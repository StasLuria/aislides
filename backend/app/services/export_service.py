"""Сервис экспорта презентаций в PDF и PPTX.

Конвертирует HTML-артефакты (слайды) в форматы PDF и PPTX.

PDF: WeasyPrint рендерит HTML → PDF (каждый слайд — отдельная страница).
PPTX: python-pptx создаёт слайды с текстовым контентом из HTML.
"""

from __future__ import annotations

import io
import logging
import re
from typing import TYPE_CHECKING

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from weasyprint import HTML

if TYPE_CHECKING:
    from backend.app.models.project import Artifact

logger = logging.getLogger(__name__)

# Размер слайда 16:9 (стандартный для презентаций)
SLIDE_WIDTH_INCHES = 13.333
SLIDE_HEIGHT_INCHES = 7.5

# CSS для PDF-рендеринга: каждый слайд на отдельной странице 16:9
PDF_WRAPPER_CSS = """
@page {
    size: 1280px 720px;
    margin: 0;
}
body {
    margin: 0;
    padding: 0;
}
.slide-page {
    width: 1280px;
    height: 720px;
    overflow: hidden;
    page-break-after: always;
    box-sizing: border-box;
}
.slide-page:last-child {
    page-break-after: auto;
}
"""


def _extract_text_blocks(html_content: str) -> list[dict[str, str]]:
    """Извлечь текстовые блоки из HTML для PPTX.

    Парсит HTML и извлекает заголовки, параграфы и списки.

    Args:
        html_content: HTML-строка слайда.

    Returns:
        Список словарей с ключами 'tag' и 'text'.
    """
    soup = BeautifulSoup(html_content, "html.parser")
    blocks: list[dict[str, str]] = []

    for element in soup.find_all(["h1", "h2", "h3", "h4", "p", "li", "span", "div"]):
        text = element.get_text(strip=True)
        if text and len(text) > 1:
            tag = element.name or "p"
            # Избегаем дублирования вложенных элементов
            if not any(b["text"] == text for b in blocks):
                blocks.append({"tag": tag, "text": text})

    return blocks


def _strip_html_tags(html_content: str) -> str:
    """Удалить HTML-теги, оставив только текст.

    Args:
        html_content: HTML-строка.

    Returns:
        Чистый текст.
    """
    return re.sub(r"<[^>]+>", "", html_content).strip()


def export_pdf(artifacts: list[Artifact]) -> bytes:
    """Экспортировать HTML-артефакты в PDF.

    Каждый артефакт (слайд) рендерится на отдельной странице 16:9.

    Args:
        artifacts: Список артефактов с HTML-контентом.

    Returns:
        PDF-файл в виде bytes.

    Raises:
        ValueError: Если нет артефактов с HTML-контентом.
    """
    html_slides = [a for a in artifacts if a.content and a.file_type == "html"]
    if not html_slides:
        msg = "Нет HTML-артефактов для экспорта в PDF"
        raise ValueError(msg)

    # Собираем все слайды в один HTML-документ
    pages = []
    for slide in html_slides:
        content = slide.content or ""
        pages.append(f'<div class="slide-page">{content}</div>')

    combined_html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <style>{PDF_WRAPPER_CSS}</style>
</head>
<body>
{''.join(pages)}
</body>
</html>"""

    logger.info("Экспорт PDF: %d слайдов", len(html_slides))

    pdf_bytes = HTML(string=combined_html).write_pdf()
    return pdf_bytes


def export_pptx(artifacts: list[Artifact]) -> bytes:
    """Экспортировать HTML-артефакты в PPTX.

    Каждый артефакт (слайд) конвертируется в текстовый слайд PPTX.
    Заголовки, параграфы и списки извлекаются из HTML.

    Args:
        artifacts: Список артефактов с HTML-контентом.

    Returns:
        PPTX-файл в виде bytes.

    Raises:
        ValueError: Если нет артефактов с HTML-контентом.
    """
    html_slides = [a for a in artifacts if a.content and a.file_type == "html"]
    if not html_slides:
        msg = "Нет HTML-артефактов для экспорта в PPTX"
        raise ValueError(msg)

    prs = Presentation()

    # Устанавливаем размер слайда 16:9
    prs.slide_width = Emu(int(SLIDE_WIDTH_INCHES * 914400))
    prs.slide_height = Emu(int(SLIDE_HEIGHT_INCHES * 914400))

    for slide_artifact in html_slides:
        content = slide_artifact.content or ""
        blocks = _extract_text_blocks(content)

        # Используем blank layout
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        if not blocks:
            # Если не удалось извлечь блоки, добавляем plain text
            plain_text = _strip_html_tags(content)
            if plain_text:
                blocks = [{"tag": "p", "text": plain_text}]

        # Размещаем текстовые блоки на слайде
        y_position = Inches(0.5)
        title_found = False

        for block in blocks:
            tag = block["tag"]
            text = block["text"]

            if tag in ("h1", "h2") and not title_found:
                # Заголовок слайда — крупный, по центру сверху
                from pptx.util import Inches as In

                txbox = slide.shapes.add_textbox(
                    In(0.5),
                    y_position,
                    In(SLIDE_WIDTH_INCHES - 1),
                    In(1.0),
                )
                tf = txbox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(32) if tag == "h1" else Pt(28)
                p.font.bold = True
                y_position += Inches(1.2)
                title_found = True

            elif tag in ("h3", "h4"):
                # Подзаголовок
                txbox = slide.shapes.add_textbox(
                    Inches(0.5),
                    y_position,
                    Inches(SLIDE_WIDTH_INCHES - 1),
                    Inches(0.6),
                )
                tf = txbox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(20)
                p.font.bold = True
                y_position += Inches(0.7)

            elif tag == "li":
                # Элемент списка
                txbox = slide.shapes.add_textbox(
                    Inches(0.8),
                    y_position,
                    Inches(SLIDE_WIDTH_INCHES - 1.3),
                    Inches(0.5),
                )
                tf = txbox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"• {text}"
                p.font.size = Pt(16)
                y_position += Inches(0.45)

            else:
                # Обычный параграф
                txbox = slide.shapes.add_textbox(
                    Inches(0.5),
                    y_position,
                    Inches(SLIDE_WIDTH_INCHES - 1),
                    Inches(0.5),
                )
                tf = txbox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(16)
                y_position += Inches(0.5)

            # Не выходим за пределы слайда
            if y_position > Inches(SLIDE_HEIGHT_INCHES - 0.5):
                break

    logger.info("Экспорт PPTX: %d слайдов", len(html_slides))

    # Сохраняем в bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
